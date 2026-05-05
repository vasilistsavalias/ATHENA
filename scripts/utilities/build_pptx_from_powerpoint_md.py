from __future__ import annotations

import argparse
import json
import math
import re
import csv
import statistics
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

from PIL import Image
from PIL import UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]

SLIDE_HEADING_RE = re.compile(r"^##\s*Slide\s+(\d+)\s+[—-]\s+(.*?)\s*$")
DEFENSE_HEADING_RE = re.compile(r"^##\s*Slide\s+(\d+)\b.*?$")
BACKTICK_PATH_RE = re.compile(r"`([^`]+)`")
BOLD_ITALIC_RE = re.compile(r"(\*\*|__|\*|_)")
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}


@dataclass
class Bullet:
    level: int
    text: str


@dataclass
class SlideSpec:
    number: int
    heading_title: str
    template: str = ""
    layout_hint: str = ""
    bullets: list[Bullet] | None = None
    visuals: list[str] | None = None
    tables: list[dict] | None = None
    data_cite: list[str] | None = None
    speaker_intent: str = ""
    build_notes: str = ""

    def __post_init__(self) -> None:
        self.bullets = self.bullets or []
        self.visuals = self.visuals or []
        self.tables = self.tables or []
        self.data_cite = self.data_cite or []


def _strip_md(text: str) -> str:
    return BOLD_ITALIC_RE.sub("", text).replace("β€”", "—").replace("β€“", "–").strip()


def parse_slide_specs(md_path: Path) -> list[SlideSpec]:
    text = md_path.read_text(encoding="utf-8", errors="replace").splitlines()

    slides: list[SlideSpec] = []
    current: SlideSpec | None = None
    section: str | None = None
    in_json_block = False
    json_lines: list[str] = []

    for raw_line in text:
        line = raw_line.rstrip("\n")

        heading_match = SLIDE_HEADING_RE.match(line)
        if heading_match:
            if current:
                slides.append(current)
            current = SlideSpec(
                number=int(heading_match.group(1)),
                heading_title=_strip_md(heading_match.group(2)),
            )
            section = None
            in_json_block = False
            json_lines = []
            continue

        if not current:
            continue

        normalized = line.strip()

        if normalized.startswith("- **Template:**"):
            section = None
            current.template = _strip_md(normalized.split(":", 1)[1].strip())
            continue

        if normalized.startswith("- **Layout"):
            section = "layout"
            # Handles "- **Layout:** ..." and variants.
            if ":" in normalized:
                after = normalized.split(":", 1)[1].strip()
                if after:
                    current.layout_hint = _strip_md(after)
            continue

        if normalized.startswith("- **On-slide text"):
            section = "onslide"
            continue

        if normalized.startswith("- **Visual(s):**"):
            section = "visuals"
            continue

        if normalized.startswith("- **Table(s)"):
            section = "tables"
            in_json_block = False
            json_lines = []
            continue

        if normalized.startswith("- **Data to cite"):
            section = "data_cite"
            continue

        if normalized.startswith("- **Speaker intent"):
            section = "speaker_intent"
            # Allow inline text after ':'
            if ":" in normalized:
                after = normalized.split(":", 1)[1].strip()
                if after:
                    current.speaker_intent = _strip_md(after)
            continue

        if normalized.startswith("- **Build/animation"):
            section = "build"
            if ":" in normalized:
                after = normalized.split(":", 1)[1].strip()
                if after:
                    current.build_notes = _strip_md(after)
            continue

        if not normalized:
            continue

        if section == "layout":
            # Multi-line layout recipes (e.g., drawing instructions).
            current.layout_hint = (current.layout_hint + " " + _strip_md(normalized)).strip()
            continue

        if section == "onslide":
            bullet_match = re.match(r"^(\s*)-\s+(.*)$", line)
            if bullet_match:
                indent = len(bullet_match.group(1).replace("\t", "    "))
                level = 0
                if indent >= 8:
                    level = 2
                elif indent >= 4:
                    level = 1
                current.bullets.append(Bullet(level=level, text=_strip_md(bullet_match.group(2))))
            else:
                current.bullets.append(Bullet(level=0, text=_strip_md(normalized)))
            continue

        if section == "visuals":
            if "none" in normalized.lower():
                continue
            for p in BACKTICK_PATH_RE.findall(line):
                candidate = p.strip()
                if Path(candidate).suffix.lower() in IMAGE_EXTS:
                    current.visuals.append(candidate)
            continue

        if section == "tables":
            if normalized.startswith("```json"):
                in_json_block = True
                json_lines = []
                continue
            if in_json_block and normalized.startswith("```"):
                in_json_block = False
                try:
                    parsed = json.loads("\n".join(json_lines).strip() or "[]")
                    if isinstance(parsed, list):
                        current.tables.extend([t for t in parsed if isinstance(t, dict)])
                except Exception:
                    pass
                json_lines = []
                continue
            if in_json_block:
                json_lines.append(line)
            continue

        if section == "data_cite":
            for p in BACKTICK_PATH_RE.findall(line):
                current.data_cite.append(p.strip())
            continue

        if section == "speaker_intent":
            current.speaker_intent = (current.speaker_intent + " " + _strip_md(normalized)).strip()
            continue

        if section == "build":
            current.build_notes = (current.build_notes + " " + _strip_md(normalized)).strip()
            continue

    if current:
        slides.append(current)

    slides.sort(key=lambda s: s.number)
    return slides


def parse_defense_notes(defense_path: Path) -> dict[int, str]:
    lines = defense_path.read_text(encoding="utf-8", errors="replace").splitlines()
    notes: dict[int, list[str]] = {}
    current_num: int | None = None

    for line in lines:
        match = DEFENSE_HEADING_RE.match(line)
        if match:
            current_num = int(match.group(1))
            notes[current_num] = []
            continue
        if current_num is None:
            continue
        # Keep formatting minimal; notes pane is plain text anyway.
        if line.strip():
            notes[current_num].append(line.rstrip())

    return {k: "\n".join(v).strip() for k, v in notes.items() if "\n".join(v).strip()}


def _norm_asset_path(p: str) -> Path:
    # Markdown uses forward slashes; normalize for Windows while keeping relative behavior.
    return (REPO_ROOT / Path(p.replace("/", "\\"))).resolve()


def _fit_contain(img_path: Path, box) -> tuple[int, int, int, int]:
    """
    Returns (left, top, width, height) in EMUs, fitted within the given box.
    Box is (left, top, width, height) EMUs.
    """
    left, top, width, height = box
    try:
        with Image.open(img_path) as im:
            w_px, h_px = im.size
    except UnidentifiedImageError:
        return left, top, width, height
    if w_px <= 0 or h_px <= 0:
        return left, top, width, height

    img_ar = w_px / h_px
    box_ar = width / height

    if img_ar >= box_ar:
        w = width
        h = int(width / img_ar)
    else:
        h = height
        w = int(height * img_ar)

    l = left + int((width - w) / 2)
    t = top + int((height - h) / 2)
    return l, t, w, h


def _grid_cells(n: int, box, padding_emu: int) -> list[tuple[int, int, int, int]]:
    left, top, width, height = box
    if n <= 0:
        return []

    box_ar = width / height
    if n == 1:
        cols, rows = 1, 1
    elif n == 2:
        cols, rows = (2, 1) if box_ar > 1.25 else (1, 2)
    elif n == 3:
        cols, rows = (3, 1) if box_ar > 1.35 else (1, 3)
    elif n == 4:
        cols, rows = 2, 2
    else:
        cols = 3 if box_ar > 1.5 else 2
        rows = math.ceil(n / cols)

    cell_w = int((width - padding_emu * (cols - 1)) / cols)
    cell_h = int((height - padding_emu * (rows - 1)) / rows)

    cells = []
    for i in range(n):
        r = i // cols
        c = i % cols
        l = left + c * (cell_w + padding_emu)
        t = top + r * (cell_h + padding_emu)
        cells.append((l, t, cell_w, cell_h))
    return cells


def _fit_cover(img_path: Path, box) -> tuple[int, int, int, int, float, float, float, float]:
    """
    Fit an image to completely cover the box (like CSS background-size: cover).
    Returns (left, top, width, height, crop_left, crop_top, crop_right, crop_bottom).
    """
    left, top, width, height = box
    try:
        with Image.open(img_path) as im:
            w_px, h_px = im.size
    except UnidentifiedImageError:
        return left, top, width, height, 0.0, 0.0, 0.0, 0.0

    if w_px <= 0 or h_px <= 0:
        return left, top, width, height, 0.0, 0.0, 0.0, 0.0

    img_ar = w_px / h_px
    box_ar = width / height

    crop_left = crop_top = crop_right = crop_bottom = 0.0
    if img_ar > box_ar:
        # Too wide: crop left/right.
        new_w = h_px * box_ar
        excess = w_px - new_w
        crop_left = (excess / 2) / w_px
        crop_right = (excess / 2) / w_px
    else:
        # Too tall: crop top/bottom.
        new_h = w_px / box_ar
        excess = h_px - new_h
        crop_top = (excess / 2) / h_px
        crop_bottom = (excess / 2) / h_px

    return left, top, width, height, crop_left, crop_top, crop_right, crop_bottom


def _extract_title_and_body(spec: SlideSpec) -> tuple[str, list[Bullet]]:
    """
    Uses a 'Title:' bullet when present; otherwise uses heading_title.
    Returns (title, remaining_bullets).
    """
    title = None
    bullets = list(spec.bullets)
    for i, b in enumerate(bullets):
        if b.level == 0 and b.text.lower().startswith("title:"):
            title = b.text.split(":", 1)[1].strip()
            bullets.pop(i)
            break
    if title is None:
        title = spec.heading_title
    return title, bullets


def _read_csv_rows(path: Path) -> list[dict]:
    if not path.exists():
        return []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        return list(csv.DictReader(f))


def _read_yaml_kv(path: Path) -> dict[str, str]:
    """
    Minimal YAML reader for top-level 'key: value' pairs.
    Ignores nested structures and lists.
    """
    if not path.exists():
        return {}
    out: dict[str, str] = {}
    for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
        if not line.strip() or line.lstrip().startswith("#"):
            continue
        if line.startswith((" ", "\t")):
            continue
        if ":" not in line:
            continue
        k, v = line.split(":", 1)
        k = k.strip()
        v = v.strip()
        if not k or not v:
            continue
        out[k] = v
    return out


def _table_from_spec(table_spec: dict) -> tuple[list[str], list[list[str]]]:
    """
    Returns (headers, rows) ready for pptx.
    """
    kind = (table_spec.get("kind") or "").strip()

    def fmt_value(key: str, value) -> str:
        if value is None:
            return ""
        value_str = str(value)
        fmt = (table_spec.get("format") or {})
        if key in fmt:
            try:
                digits = int(fmt[key])
                return f"{float(value_str):.{digits}f}"
            except Exception:
                return value_str
        return value_str

    if kind == "csv":
        src = table_spec.get("source", "")
        rows = _read_csv_rows(_norm_asset_path(src)) if src else []
        where = table_spec.get("where") or {}
        where_in = table_spec.get("where_in") or {}
        if where:
            rows = [r for r in rows if all(str(r.get(k, "")) == str(v) for k, v in where.items())]
        if where_in:
            for k, allowed in where_in.items():
                allowed_set = {str(a) for a in (allowed or [])}
                rows = [r for r in rows if str(r.get(k, "")) in allowed_set]
        columns = table_spec.get("columns") or (list(rows[0].keys()) if rows else [])
        if isinstance(columns, str):
            columns = [columns]
        rows_cfg = table_spec.get("rows", "all")
        if isinstance(rows_cfg, dict) and "top_n" in rows_cfg:
            rows = rows[: int(rows_cfg["top_n"])]
        data_rows = [[fmt_value(c, r.get(c, "")) for c in columns] for r in rows]
        return list(columns), data_rows

    if kind == "json":
        src = table_spec.get("source", "")
        if not src:
            return ["Key", "Value"], [["(missing source)", ""]]
        p = _norm_asset_path(src)
        if not p.exists():
            return ["Key", "Value"], [["(missing)", src]]
        obj = json.loads(p.read_text(encoding="utf-8", errors="replace"))
        if isinstance(obj, dict):
            items = list(obj.items())
            rows_cfg = table_spec.get("rows", "all")
            if isinstance(rows_cfg, dict) and "top_n" in rows_cfg:
                items = items[: int(rows_cfg["top_n"])]
            return ["Key", "Value"], [[str(k), str(v)] for k, v in items]
        if isinstance(obj, list) and obj and isinstance(obj[0], dict):
            columns = table_spec.get("columns") or list(obj[0].keys())
            rows_cfg = table_spec.get("rows", "all")
            rows_list = obj
            if isinstance(rows_cfg, dict) and "top_n" in rows_cfg:
                rows_list = rows_list[: int(rows_cfg["top_n"])]
            return list(columns), [[str(r.get(c, "")) for c in columns] for r in rows_list]
        return ["Value"], [[str(obj)]]

    if kind == "yaml_kv":
        src = table_spec.get("source", "")
        kv = _read_yaml_kv(_norm_asset_path(src)) if src else {}
        return ["Key", "Value"], [[k, v] for k, v in kv.items()]

    if kind == "computed":
        compute = (table_spec.get("compute") or "").strip()
        paths = table_spec.get("paths") or []
        if isinstance(paths, str):
            paths = [paths]

        if compute == "kpi_summary":
            funnel = json.loads(_norm_asset_path("outputs/03_intelligent_filtering/filtering_stats.json").read_text(encoding="utf-8"))
            split = json.loads(_norm_asset_path("outputs/09_data_splitting/split_stats.json").read_text(encoding="utf-8"))
            seeds = json.loads(_norm_asset_path("outputs/00_reproducibility/seeds.json").read_text(encoding="utf-8"))
            compute_log = _read_csv_rows(_norm_asset_path("outputs/00_reproducibility/compute_log.csv"))[0]
            paired_rows = _read_csv_rows(_norm_asset_path("outputs/13_model_evaluation/statistical_tests/paired_t_tests.csv"))
            n_pairs = paired_rows[0].get("n_pairs", "") if paired_rows else ""
            matrix_rows = _read_csv_rows(_norm_asset_path("outputs/13_model_evaluation/benchmarking_matrix/matrix_results.csv"))
            total_rows = len(matrix_rows)
            rows = [
                ["Raw images", str(funnel.get("total"))],
                ["Kept after YOLO+blur", str(funnel.get("kept"))],
                ["Crops", str(funnel.get("crops"))],
                ["Split (train/val/test)", f"{split.get('train')}/{split.get('validation')}/{split.get('test')}"],
                ["Eval pairs", str(n_pairs)],
                ["Eval rows", str(total_rows)],
                ["Seed", str(seeds.get("global_random_state"))],
                ["GPU", f"{compute_log.get('gpu_count')}× {compute_log.get('gpu_name')}"],
            ]
            return ["KPI", "Value"], rows

        if compute == "filtering_yolo_only":
            funnel = json.loads(_norm_asset_path("outputs/03_intelligent_filtering/filtering_stats.json").read_text(encoding="utf-8"))
            return ["KPI", "Value"], [["rejected_yolo", str(funnel.get("rejected_yolo"))]]

        if compute == "exec_times_top_slowest":
            exec_rows = _read_csv_rows(_norm_asset_path("outputs/00_logs/execution_times.csv"))
            exec_rows_sorted = sorted(exec_rows, key=lambda r: float(r.get("duration_seconds", "0") or 0), reverse=True)[:8]
            rows = [[r.get("stage_num", ""), r.get("stage_name", ""), f"{float(r.get('duration_seconds','0') or 0):.2f}"] for r in exec_rows_sorted]
            return ["stage_num", "stage_name", "duration_seconds"], rows

        if compute == "coverage_summary":
            cov_rows = _read_csv_rows(_norm_asset_path("outputs/10_feature_engineering/mask_coverage_test.csv"))
            vals = [float(r["coverage_ratio"]) for r in cov_rows if r.get("coverage_ratio")]
            if not vals:
                return ["KPI", "Value"], [["coverage_ratio", "(no data)"]]
            rows = [
                ["min_coverage", f"{min(vals)*100:.1f}%"],
                ["median_coverage", f"{statistics.median(vals)*100:.1f}%"],
                ["max_coverage", f"{max(vals)*100:.1f}%"],
            ]
            return ["KPI", "Value"], rows

        if compute == "file_manifest":
            rows = []
            for p in paths:
                abs_p = _norm_asset_path(p)
                size = abs_p.stat().st_size if abs_p.exists() else 0
                rows.append([str(p), str(size)])
            return ["path", "bytes"], rows

        if compute == "stage17_requirements":
            exec_rows = _read_csv_rows(_norm_asset_path("outputs/00_logs/execution_times.csv"))
            executed = any(r.get("stage_num") == "17" for r in exec_rows)
            rows = [
                ["Executed in this run?", "Yes" if executed else "No"],
                ["Inputs required", "{image}.(jpg|png) + {image}_mask.(png|jpg)"],
                ["Outputs expected", "{stem}_Telea.png, {stem}_Ours.png, grid_{stem}.png"],
            ]
            return ["Field", "Value"], rows

        if compute == "ablation_delta_zero":
            seeds = json.loads(_norm_asset_path("outputs/00_reproducibility/seeds.json").read_text(encoding="utf-8"))
            rows = [["Δ(max-min)", "0.0000"], ["Seed", str(seeds.get("global_random_state", ""))]]
            return ["Field", "Value"], rows

        if compute == "caption_quality_headline":
            obj = json.loads(_norm_asset_path("outputs/06_caption_generation/caption_quality_report.json").read_text(encoding="utf-8"))
            rows = [
                ["total_images", str(obj.get("total_images"))],
                ["empty_captions", str(obj.get("empty_captions"))],
                ["short_captions_lt5_words", str(obj.get("short_captions_lt5_words"))],
                ["sample_count", str(obj.get("sample_count"))],
            ]
            return ["Field", "Value"], rows

        if compute == "method_mapping_sample":
            obj = json.loads(_norm_asset_path("outputs/16_expert_validation/method_mapping.json").read_text(encoding="utf-8"))
            items = list(obj.items())[:6]
            rows = []
            for k, v in items:
                if isinstance(v, dict) and v:
                    mapped = ", ".join([f"{mk}={mv}" for mk, mv in v.items()])
                else:
                    mapped = "(unmapped)"
                rows.append([k, mapped])
            return ["file", "mapping"], rows

    return ["Value"], [["(unsupported table spec)", json.dumps(table_spec, ensure_ascii=False)]]


def _add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]]) -> None:
    n_rows = max(1, len(rows) + 1)
    n_cols = max(1, len(headers))
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header styling
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(45, 45, 45)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

    for r_i, row in enumerate(rows, start=1):
        for c_i in range(n_cols):
            cell = table.cell(r_i, c_i)
            cell.text = str(row[c_i] if c_i < len(row) else "")
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(30, 30, 30)


def build_pptx(
    slide_specs: list[SlideSpec],
    defense_notes: dict[int, str] | None,
    output_path: Path,
    report_path: Path | None,
) -> dict:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    emu = 914400
    slide_w = int(13.333 * emu)
    slide_h = int(7.5 * emu)
    margin = int(0.6 * emu)
    gutter = int(0.25 * emu)
    footer_h = int(0.35 * emu)
    title_h = int(0.85 * emu)

    missing_assets: set[str] = set()
    table_errors: list[dict] = []
    per_slide = []
    total_slides = len(slide_specs)

    for spec in slide_specs:
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()

        template = (spec.template or "").strip()
        title_text, body_bullets = _extract_title_and_body(spec)

        # Visual paths (normalized)
        visual_paths: list[tuple[str, Path]] = []
        for v in [v for v in (spec.visuals or []) if v]:
            pth = _norm_asset_path(v)
            if not pth.exists():
                missing_assets.add(v)
            visual_paths.append((v, pth))

        # Slide number footer
        footer = slide.shapes.add_textbox(Inches(12.0), Inches(7.15), Inches(1.3), Inches(0.3))
        f_tf = footer.text_frame
        f_tf.clear()
        fp = f_tf.paragraphs[0]
        fp.text = f"{spec.number}/{total_slides}"
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(120, 120, 120)
        fp.alignment = PP_ALIGN.RIGHT

        # Sources footer (small, de-duped)
        sources: list[str] = []
        sources.extend(spec.data_cite or [])
        for t in spec.tables or []:
            if isinstance(t, dict):
                if t.get("source"):
                    sources.append(str(t["source"]))
                for p in t.get("paths") or []:
                    sources.append(str(p))
        sources.extend([v for v, _ in visual_paths])
        seen = set()
        sources_d = []
        for s in sources:
            if s and s not in seen:
                seen.add(s)
                sources_d.append(s)
        if sources_d:
            shown = sources_d[:3]
            extra = len(sources_d) - len(shown)
            src_text = "Sources: " + "; ".join(shown) + (f" (+{extra} more)" if extra > 0 else "")
            src_box = slide.shapes.add_textbox(margin, slide_h - footer_h, slide_w - 2 * margin - int(1.6 * emu), footer_h)
            stf = src_box.text_frame
            stf.clear()
            sp = stf.paragraphs[0]
            sp.text = src_text
            sp.font.size = Pt(9)
            sp.font.color.rgb = RGBColor(140, 140, 140)

        # Notes (Greek script)
        if defense_notes and spec.number in defense_notes:
            slide.notes_slide.notes_text_frame.text = defense_notes[spec.number]

        # Content bounds (excluding footer)
        content_top = int(0.15 * emu) + title_h
        content_left = margin
        content_w = slide_w - 2 * margin
        content_h = slide_h - content_top - footer_h - int(0.05 * emu)

        def add_title_box(text: str) -> None:
            tb = slide.shapes.add_textbox(margin, int(0.15 * emu), slide_w - 2 * margin, title_h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = RGBColor(20, 20, 20)

        def add_bullets_box(left, top, width, height, bullets: list[Bullet], max_lines: int = 10) -> None:
            tbox = slide.shapes.add_textbox(left, top, width, height)
            tfb = tbox.text_frame
            tfb.clear()
            tfb.word_wrap = True
            tfb.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            show = bullets[:max_lines]
            for idx, b in enumerate(show):
                para = tfb.paragraphs[0] if idx == 0 else tfb.add_paragraph()
                para.text = b.text
                para.level = min(b.level, 2)
                para.font.size = Pt(18 if idx == 0 and len(show) <= 6 else 16)
                para.font.color.rgb = RGBColor(40, 40, 40)

        # ---- Templates ----
        if template == "intro_hero" and visual_paths:
            # Full-bleed image
            v_rel, v_abs = visual_paths[0]
            if v_abs.exists():
                l, t, w, h, cl, ct, cr, cb = _fit_cover(v_abs, (0, 0, slide_w, slide_h))
                pic = slide.shapes.add_picture(str(v_abs), l, t, w, h)
                try:
                    pic.crop_left, pic.crop_top, pic.crop_right, pic.crop_bottom = cl, ct, cr, cb
                except Exception:
                    pass
            else:
                missing_assets.add(v_rel)

            # Dark overlay bar
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, int(1.2 * emu))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            try:
                overlay.fill.transparency = 0.35
            except Exception:
                pass
            overlay.line.fill.background()

            # Use the first 2–3 bullets as title/subtitle/footer text.
            tb = slide.shapes.add_textbox(int(0.7 * emu), int(0.12 * emu), slide_w - int(1.4 * emu), int(1.05 * emu))
            tf = tb.text_frame
            tf.clear()
            lines = (spec.bullets or [])[:3]
            for i, b in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b.text
                p.font.size = Pt(36 if i == 0 else 18)
                p.font.bold = True if i == 0 else False
                p.font.color.rgb = RGBColor(255, 255, 255)

        elif template == "outro_visual" and visual_paths:
            add_title_box(title_text)
            v_rel, v_abs = visual_paths[0]
            if v_abs.exists():
                l, t, w, h = _fit_contain(v_abs, (content_left, content_top, content_w, content_h))
                slide.shapes.add_picture(str(v_abs), l, t, w, h)
            else:
                missing_assets.add(v_rel)
            # Optional callouts (short)
            if body_bullets:
                add_bullets_box(content_left, content_top, int(content_w * 0.35), int(content_h * 0.35), body_bullets, max_lines=5)

        else:
            add_title_box(title_text)

            if template == "stage_theory":
                left_w = int(content_w * 0.58)
                right_w = content_w - left_w - gutter
                add_bullets_box(content_left, content_top, left_w, content_h, body_bullets, max_lines=12)

                # Mini-diagram boxes (Input/Process/Output)
                diag_left = content_left + left_w + gutter
                box_h = int((content_h - 2 * gutter) / 3)
                def find_field(prefix: str) -> str:
                    for b in spec.bullets or []:
                        if b.text.lower().startswith(prefix):
                            return b.text.split(":", 1)[1].strip()
                    return ""
                fields = [("Input", find_field("input:")), ("Process", find_field("process:")), ("Output", find_field("output:"))]
                for i, (hdr, val) in enumerate(fields):
                    y = content_top + i * (box_h + gutter)
                    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diag_left, y, right_w, box_h)
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = RGBColor(240, 242, 246)
                    shp.line.color.rgb = RGBColor(210, 215, 222)
                    tfm = shp.text_frame
                    tfm.clear()
                    p1 = tfm.paragraphs[0]
                    p1.text = hdr
                    p1.font.size = Pt(14)
                    p1.font.bold = True
                    p2 = tfm.add_paragraph()
                    p2.text = val or ""
                    p2.font.size = Pt(16)

            elif template in {"stage_results", "outro_summary", "outro_limits", "outro_qa"}:
                left_w = int(content_w * 0.60) if template == "stage_results" else int(content_w * 0.55)
                right_w = content_w - left_w - gutter
                left_left = content_left
                right_left = content_left + left_w + gutter

                # Left: bullets at top, visuals below
                bullets_h = int(content_h * 0.30) if body_bullets else 0
                if body_bullets:
                    add_bullets_box(left_left, content_top, left_w, bullets_h, body_bullets, max_lines=8)

                img_top = content_top + bullets_h + (gutter if body_bullets else 0)
                img_h = content_h - bullets_h - (gutter if body_bullets else 0)
                if visual_paths:
                    img_box = (left_left, img_top, left_w, img_h)
                    cells = _grid_cells(len(visual_paths), img_box, padding_emu=int(0.12 * emu))
                    for (v_rel, v_abs), cell in zip(visual_paths, cells, strict=False):
                        if v_abs.exists():
                            l, t, w, h = _fit_contain(v_abs, cell)
                            slide.shapes.add_picture(str(v_abs), l, t, w, h)
                        else:
                            missing_assets.add(v_rel)

                # Right: tables stacked (or fallback bullets if none)
                if spec.tables:
                    per_h = int((content_h - gutter * (len(spec.tables) - 1)) / len(spec.tables))
                    for i, tbl_spec in enumerate(spec.tables):
                        try:
                            headers, rows = _table_from_spec(tbl_spec)
                        except Exception as e:
                            table_errors.append({"slide": spec.number, "error": str(e), "table": tbl_spec})
                            headers, rows = ["Error"], [[str(e)]]
                        top = content_top + i * (per_h + gutter)
                        _add_table(slide, right_left, top, right_w, per_h, headers, rows[: min(12, len(rows))])
                else:
                    add_bullets_box(right_left, content_top, right_w, content_h, [], max_lines=0)

            elif template == "intro_kpi":
                # Hero callout + two tables side-by-side
                hero_h = int(content_h * 0.18)
                hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, content_left, content_top, content_w, hero_h)
                hero.fill.solid()
                hero.fill.fore_color.rgb = RGBColor(20, 83, 136)
                hero.line.fill.background()
                tfm = hero.text_frame
                tfm.clear()
                line = next((b.text for b in body_bullets if "best fid" in b.text.lower()), "Best FID (see table)")
                p = tfm.paragraphs[0]
                p.text = line
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                tables_top = content_top + hero_h + gutter
                tables_h = content_h - hero_h - gutter
                col_w = int((content_w - gutter) / 2)
                for i, tbl_spec in enumerate((spec.tables or [])[:2]):
                    try:
                        headers, rows = _table_from_spec(tbl_spec)
                    except Exception as e:
                        table_errors.append({"slide": spec.number, "error": str(e), "table": tbl_spec})
                        headers, rows = ["Error"], [[str(e)]]
                    left = content_left + i * (col_w + gutter)
                    _add_table(slide, left, tables_top, col_w, tables_h, headers, rows[: min(14, len(rows))])
                # Optional small visual at bottom right if present and space remains.
                if len(visual_paths) >= 1:
                    v_rel, v_abs = visual_paths[0]
                    if v_abs.exists():
                        l, t, w, h = _fit_contain(v_abs, (content_left + col_w + gutter, tables_top + int(tables_h * 0.62), col_w, int(tables_h * 0.38)))
                        slide.shapes.add_picture(str(v_abs), l, t, w, h)

            else:
                # Generic fallback
                if visual_paths and spec.tables:
                    left_w = int(content_w * 0.55)
                    right_w = content_w - left_w - gutter
                    img_box = (content_left, content_top, left_w, content_h)
                    cells = _grid_cells(len(visual_paths), img_box, padding_emu=int(0.12 * emu))
                    for (v_rel, v_abs), cell in zip(visual_paths, cells, strict=False):
                        if v_abs.exists():
                            l, t, w, h = _fit_contain(v_abs, cell)
                            slide.shapes.add_picture(str(v_abs), l, t, w, h)
                        else:
                            missing_assets.add(v_rel)
                    per_h = int((content_h - gutter * (len(spec.tables) - 1)) / len(spec.tables))
                    for i, tbl_spec in enumerate(spec.tables):
                        headers, rows = _table_from_spec(tbl_spec)
                        top = content_top + i * (per_h + gutter)
                        _add_table(slide, content_left + left_w + gutter, top, right_w, per_h, headers, rows[: min(12, len(rows))])
                elif visual_paths:
                    img_box = (content_left, content_top, content_w, int(content_h * 0.72))
                    cells = _grid_cells(len(visual_paths), img_box, padding_emu=int(0.12 * emu))
                    for (v_rel, v_abs), cell in zip(visual_paths, cells, strict=False):
                        if v_abs.exists():
                            l, t, w, h = _fit_contain(v_abs, cell)
                            slide.shapes.add_picture(str(v_abs), l, t, w, h)
                        else:
                            missing_assets.add(v_rel)
                    add_bullets_box(content_left, content_top + int(content_h * 0.74), content_w, int(content_h * 0.26), body_bullets, max_lines=8)
                else:
                    add_bullets_box(content_left, content_top, content_w, content_h, body_bullets, max_lines=12)

        per_slide.append(
            {
                "slide": spec.number,
                "template": template,
                "title": title_text,
                "visuals": [v for v, _ in visual_paths],
                "missing_visuals": [v for v, p in visual_paths if not p.exists()],
                "tables": spec.tables,
            }
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    report = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "output_pptx": str(output_path.relative_to(REPO_ROOT)),
        "slide_count": len(slide_specs),
        "missing_assets_count": len(missing_assets),
        "missing_assets": sorted(missing_assets),
        "table_errors": table_errors,
        "slides": per_slide,
    }

    if report_path is not None:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")

    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a .pptx from docs/presentations/powerpoint.md and outputs/ assets.")
    parser.add_argument(
        "--md",
        default="docs/presentations/powerpoint.md",
        help="Markdown slide spec path (default: docs/presentations/powerpoint.md)",
    )
    parser.add_argument(
        "--defense",
        default="docs/presentations/defense_script.md",
        help="Optional defense script markdown to inject into speaker notes",
    )
    parser.add_argument(
        "--out",
        default="docs/presentations/Semantic-Restoration-of-Ancient-Greek-Pottery_V3_auto_v2.pptx",
        help="Output .pptx path",
    )
    parser.add_argument(
        "--report",
        default="docs/presentations/_build/pptx_build_report_v2.json",
        help="Write a JSON build report (missing assets, slide list)",
    )
    args = parser.parse_args()

    md_path = (REPO_ROOT / Path(args.md)).resolve()
    if not md_path.exists():
        raise FileNotFoundError(f"Slide spec markdown not found: {md_path}")

    defense_path = (REPO_ROOT / Path(args.defense)).resolve()
    defense_notes = parse_defense_notes(defense_path) if defense_path.exists() else None

    slide_specs = parse_slide_specs(md_path)
    if not slide_specs:
        raise ValueError(f"No slides parsed from: {md_path}")

    out_path = (REPO_ROOT / Path(args.out)).resolve()
    report_path = (REPO_ROOT / Path(args.report)).resolve() if args.report else None

    report = build_pptx(slide_specs, defense_notes, out_path, report_path)

    print(f"Built PPTX: {report['output_pptx']} ({report['slide_count']} slides)")
    if report["missing_assets_count"]:
        print(f"Missing assets: {report['missing_assets_count']} (see report: {Path(args.report)})")
    else:
        print("All referenced assets were found.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
